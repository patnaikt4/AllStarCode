"""
Generate and upsert curriculum embeddings for the AI Lesson Video Feedback Tool.

Reads curriculum documents (PPTX, PDF, TXT, or Markdown), splits them into
chunks, generates embeddings via OpenAI, and stores them in Supabase for
later similarity search.

Supports nested folder structures like:
    curriculum/
        Week 1/
            Day 1/
                01 What is CS_/
                    Lecture_ What is CS_.pptx

Usage:
    python generate_embeddings.py --docs-dir ./curriculum_docs
    python generate_embeddings.py --file ./curriculum_docs/week3.pptx
"""

import argparse
import os
import sys
from pathlib import Path
from typing import Dict, List

from dotenv import load_dotenv
from openai import OpenAI
from supabase import create_client

# ── Load environment variables ───────────────────────────────────────────
# Same .env.local path convention used by similarity_search.py
load_dotenv(
    dotenv_path=os.path.join(
        os.path.dirname(os.path.abspath(__file__)),
        "..",
        "apps",
        "web",
        ".env.local",
    )
)

# ── Constants ────────────────────────────────────────────────────────────
EMBEDDING_MODEL = "text-embedding-3-small"
EMBEDDING_DIMENSIONS = 1536
# How many characters per chunk (with overlap so we don't split mid-thought)
CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200

# ── Initialise API clients ───────────────────────────────────────────────
_openai = OpenAI(api_key=os.environ["OPENAI_API_KEY"])
_supabase = create_client(
    os.environ["NEXT_PUBLIC_SUPABASE_URL"],
    os.environ["NEXT_PUBLIC_SUPABASE_ANON_KEY"],
)


# ── Core functions ───────────────────────────────────────────────────────


def embed_text(text: str) -> List[float]:
    """
    Generate embedding vector for given text using OpenAI.

    Sends the text to OpenAI's text-embedding-3-small model and returns
    a list of 1536 floating-point numbers that represent the text's
    "meaning" in vector space.

    Args:
        text (str): Input text to embed.

    Returns:
        List[float]: Embedding vector (1536 dimensions).
    """
    response = _openai.embeddings.create(
        model=EMBEDDING_MODEL,
        input=text,
        dimensions=EMBEDDING_DIMENSIONS,
    )
    return response.data[0].embedding


def upsert_curriculum_chunk(
    source_doc: str,
    chunk_text: str,
    embedding: List[float],
    metadata: Dict,
) -> None:
    """
    Store a curriculum chunk and its embedding in the Supabase
    curriculum_chunks table.

    Uses upsert logic: if a row with the same source_doc + chunk_text
    already exists, it will be inserted as a new row (the table uses
    auto-generated UUIDs, so duplicates are possible — see the note
    about deduplication in main()).

    Args:
        source_doc (str): Name of the source document (e.g. "Week 3 - Loops.pdf").
        chunk_text (str): The actual text content of this chunk.
        embedding (List[float]): The 1536-dim embedding vector.
        metadata (Dict): Extra info like {"section": "Week 3", "chunk_index": 0}.
    """
    _supabase.table("curriculum_chunks").insert(
        {
            "source_doc": source_doc,
            "chunk_text": chunk_text,
            "embedding": embedding,
            "metadata": metadata,
        }
    ).execute()


# ── Text processing helpers ──────────────────────────────────────────────


def read_file_text(file_path: Path) -> str:
    """
    Read text content from a file. Supports .pptx, .pdf, .txt, and .md.

    For PPTX, uses python-pptx to extract text from every slide's shapes.
    For PDFs, uses PyPDF2 to extract text page by page.
    For text/markdown, reads the file directly as UTF-8.

    Args:
        file_path (Path): Path to the file to read.

    Returns:
        str: The extracted text content.
    """
    suffix = file_path.suffix.lower()

    if suffix == ".pptx":
        try:
            from pptx import Presentation
        except ImportError:
            print(
                "python-pptx is required for PPTX files. Install it with: "
                "pip install python-pptx",
                file=sys.stderr,
            )
            sys.exit(1)

        prs = Presentation(str(file_path))
        slides_text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_parts = [f"[Slide {slide_num}]"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_parts.append(text)
                # Also extract text from tables in slides
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_text:
                            slide_parts.append(" | ".join(row_text))
            slides_text.append("\n".join(slide_parts))
        return "\n\n".join(slides_text)

    if suffix == ".pdf":
        try:
            from PyPDF2 import PdfReader
        except ImportError:
            print(
                "PyPDF2 is required for PDF files. Install it with: "
                "pip install PyPDF2",
                file=sys.stderr,
            )
            sys.exit(1)

        reader = PdfReader(str(file_path))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n\n".join(pages)

    # .txt, .md, or any other text format
    return file_path.read_text(encoding="utf-8")


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    """
    Split a long text into overlapping chunks.

    Why overlap? If we split exactly at position 1000, we might cut a
    sentence in half. The overlap means each chunk shares some text with
    the next one, so important context near the boundary isn't lost.

    Example with chunk_size=10, overlap=3:
        "Hello world, how are you today?"
        Chunk 1: "Hello worl"
        Chunk 2: "orl, how a"   (starts 3 characters before the end of chunk 1)
        Chunk 3: "ow are you"
        ...

    Args:
        text (str): The full text to split.
        chunk_size (int): Max characters per chunk (default 1000).
        overlap (int): How many characters to share between consecutive chunks.

    Returns:
        List[str]: List of text chunks.
    """
    # Clean up excessive whitespace
    text = text.strip()
    if not text:
        return []

    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        # Move forward by (chunk_size - overlap) so chunks overlap
        start += chunk_size - overlap

    return chunks


def clear_existing_chunks(source_doc: str) -> int:
    """
    Delete all existing chunks for a given source document.

    This prevents duplicates when re-running the script on the same
    document (e.g. after curriculum updates).

    Args:
        source_doc (str): The document name to clear.

    Returns:
        int: Number of rows deleted.
    """
    result = (
        _supabase.table("curriculum_chunks")
        .delete()
        .eq("source_doc", source_doc)
        .execute()
    )
    return len(result.data) if result.data else 0


# ── Main entry point ─────────────────────────────────────────────────────


def process_file(file_path: Path, base_dir: Path | None = None) -> int:
    """
    Process a single curriculum file: read → chunk → embed → store.

    Args:
        file_path (Path): Path to the curriculum document.
        base_dir (Path | None): If provided, source_doc will be the relative
            path from base_dir (e.g. "Week 1/Day 1/Lecture.pptx").

    Returns:
        int: Number of chunks embedded and stored.
    """
    # Use relative path as source_doc so we preserve folder context
    # e.g. "Week 1/Day 1/01 What is CS_/Lecture_ What is CS_.pptx"
    if base_dir:
        source_doc = str(file_path.relative_to(base_dir))
    else:
        source_doc = file_path.name
    print(f"\n📄 Processing: {source_doc}")

    # Step 1: Read the file
    text = read_file_text(file_path)
    if not text.strip():
        print(f"   ⚠️  Skipping {source_doc} — file is empty.")
        return 0

    # Step 2: Split into chunks
    chunks = chunk_text(text)
    print(f"   Split into {len(chunks)} chunks (chunk_size={CHUNK_SIZE}, overlap={CHUNK_OVERLAP})")

    # Step 3: Clear old embeddings for this doc (prevents duplicates on re-run)
    deleted = clear_existing_chunks(source_doc)
    if deleted > 0:
        print(f"   Cleared {deleted} existing chunks for this document")

    # Step 4: Embed each chunk and store it
    for i, chunk in enumerate(chunks):
        print(f"   Embedding chunk {i + 1}/{len(chunks)}...", end=" ")
        embedding = embed_text(chunk)
        metadata = {
            "source": source_doc,
            "chunk_index": i,
            "total_chunks": len(chunks),
        }
        upsert_curriculum_chunk(source_doc, chunk, embedding, metadata)
        print("✓")

    print(f"   ✅ Done — {len(chunks)} chunks stored for {source_doc}")
    return len(chunks)


def main() -> int:
    """
    CLI entry point for generating curriculum embeddings.

    Supports two modes:
      --file <path>       Process a single file
      --docs-dir <path>   Process all .txt, .md, and .pdf files in a directory
    """
    parser = argparse.ArgumentParser(
        description="Generate and store curriculum embeddings in Supabase."
    )
    parser.add_argument(
        "--file",
        type=str,
        help="Path to a single curriculum document (.pptx, .txt, .md, or .pdf).",
    )
    parser.add_argument(
        "--docs-dir",
        type=str,
        help="Path to a directory of curriculum documents (searched recursively).",
    )
    parser.add_argument(
        "--chunk-size",
        type=int,
        default=CHUNK_SIZE,
        help=f"Characters per chunk (default: {CHUNK_SIZE}).",
    )
    parser.add_argument(
        "--chunk-overlap",
        type=int,
        default=CHUNK_OVERLAP,
        help=f"Overlap between chunks (default: {CHUNK_OVERLAP}).",
    )
    args = parser.parse_args()

    # Update globals if custom values provided
    global CHUNK_SIZE, CHUNK_OVERLAP
    CHUNK_SIZE = args.chunk_size
    CHUNK_OVERLAP = args.chunk_overlap

    if not args.file and not args.docs_dir:
        parser.error("Provide either --file or --docs-dir.")

    # Collect files to process
    files: List[Path] = []

    if args.file:
        p = Path(args.file)
        if not p.exists():
            print(f"Error: File not found: {p}", file=sys.stderr)
            return 1
        files.append(p)

    if args.docs_dir:
        d = Path(args.docs_dir)
        if not d.is_dir():
            print(f"Error: Directory not found: {d}", file=sys.stderr)
            return 1
        # Use ** glob pattern to search recursively through nested folders
        # (e.g. Week 1/Day 1/01 What is CS_/Lecture.pptx)
        for ext in ("**/*.pptx", "**/*.txt", "**/*.md", "**/*.pdf"):
            files.extend(sorted(d.glob(ext)))

    if not files:
        print("No supported files found (.pptx, .txt, .md, .pdf).", file=sys.stderr)
        return 1

    # Determine base_dir for relative path naming
    base_dir = Path(args.docs_dir) if args.docs_dir else None

    print(f"Found {len(files)} file(s) to process.")

    total_chunks = 0
    for file_path in files:
        try:
            total_chunks += process_file(file_path, base_dir=base_dir)
        except Exception as exc:
            print(f"   ❌ Error processing {file_path.name}: {exc}", file=sys.stderr)
            continue

    print(f"\n🎉 All done! Embedded {total_chunks} total chunks from {len(files)} file(s).")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())